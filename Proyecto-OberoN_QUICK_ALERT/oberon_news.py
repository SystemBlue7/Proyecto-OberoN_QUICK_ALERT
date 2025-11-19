import os
import locale
import win32com.client
import tkinter as tk
from tkinter import ttk, filedialog, scrolledtext, messagebox
from pptx import Presentation
from pptx.util import Cm, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from PIL import Image, ImageTk
from datetime import datetime

# ============================================================
# CONFIGURACIÓN DE IDIOMA
# ============================================================
for loc in ("es_ES.UTF-8", "es_CO.UTF-8", "es_ES.utf8", "Spanish_Colombia", "Spanish_Spain"):
    try:
        locale.setlocale(locale.LC_TIME, loc)
        break
    except Exception:
        continue

CARPETA_CABEZOTE = os.path.join(os.path.dirname(__file__), "OBERON NEWS.png")
LOGO_EMPRESA = os.path.join(os.path.dirname(__file__), "Logo.png")

# ============================================================
# FUNCIÓN DE CREACIÓN DEL BOLETÍN
# ============================================================
def generar_boletin():
    fecha = entry_fecha.get().strip()
    titulo = entry_titulo.get().strip()
    imagen_evidencia = entry_imagen.get().strip()
    archivo_salida = entry_nombre.get().strip()

    if not all([fecha, titulo, imagen_evidencia, archivo_salida]):
        messagebox.showerror("Error", "Por favor diligencia todos los campos antes de generar el boletín.")
        return

    modelo = Presentation()
    modelo.slide_width = Cm(38.1)
    modelo.slide_height = Cm(67.733)
    slide = modelo.slides.add_slide(modelo.slide_layouts[6])

    def agregar_texto_avanzado(slide, x, y, ancho, alto, widget, tamano, color_base, centrado=False):
        caja = slide.shapes.add_textbox(Cm(x), Cm(y), Cm(ancho), Cm(alto))
        tf = caja.text_frame
        tf.word_wrap = True

        if isinstance(widget, tk.Entry):
            contenido = widget.get().strip()
        else:
            contenido = widget.get("1.0", tk.END).strip()

        if not contenido:
            return

        lineas = contenido.split("\n")
        for i, linea in enumerate(lineas):
            p = tf.add_paragraph()
            p.font.name = "Century Gothic"
            p.font.size = Pt(tamano)
            p.alignment = PP_ALIGN.CENTER if centrado else PP_ALIGN.LEFT

            palabras = linea.split(" ")
            index_offset = 0
            for palabra in palabras:
                run = p.add_run()
                run.text = palabra + " "
                index_start = f"{i+1}.{index_offset}"
                index_offset += len(palabra) + 1

                tags = []
                if not isinstance(widget, tk.Entry):
                    for tag in widget.tag_names():
                        try:
                            if widget.tag_ranges(tag):
                                for j in range(0, len(widget.tag_ranges(tag)), 2):
                                    start = widget.tag_ranges(tag)[j]
                                    end = widget.tag_ranges(tag)[j + 1]
                                    if widget.compare(index_start, ">=", start) and widget.compare(index_start, "<", end):
                                        tags.append(tag)
                        except tk.TclError:
                            pass

                run.font.name = "Century Gothic"
                run.font.size = Pt(tamano)
                run.font.color.rgb = color_base
                run.font.bold = "bold" in tags
                if "rojo" in tags:
                    run.font.color.rgb = RGBColor(214, 65, 35)
                elif "azul" in tags:
                    run.font.color.rgb = RGBColor(0, 48, 87)

    # ======================= CONTENIDO SLIDE =======================
    slide.shapes.add_picture(CARPETA_CABEZOTE, Cm(0), Cm(0), Cm(38.1), Cm(9.04))
    agregar_texto_avanzado(slide, 7.89, 5.03, 14.14, 1.51, entry_fecha, 30, RGBColor(0, 0, 0), centrado=True)
    agregar_texto_avanzado(slide, 2.02, 7.58, 18.59, 3.3, entry_titulo, 30, RGBColor(0, 48, 87))

    def agregar_linea(slide, y):
        caja = slide.shapes.add_textbox(Cm(1.81), Cm(y), Cm(33.84), Cm(0))
        tf = caja.text_frame
        p = tf.paragraphs[0]
        run = p.add_run()
        run.text = "───────────────────────────────────────────────"
        run.font.name = "Century Gothic"
        run.font.size = Pt(28)
        run.font.color.rgb = RGBColor(214, 65, 35)
        p.alignment = PP_ALIGN.CENTER

    agregar_texto_avanzado(slide, 2.17, 11.07, 33.84, 8.34, txt_texto1, 24, RGBColor(0, 48, 87))
    agregar_linea(slide, 20.80)
    agregar_texto_avanzado(slide, 1.87, 21.56, 16.61, 14.12, txt_texto2, 22, RGBColor(0, 48, 87))
    slide.shapes.add_picture(imagen_evidencia, Cm(20.2), Cm(22.79), Cm(15.45), Cm(11.7))
    agregar_texto_avanzado(slide, 2.02, 35.25, 33.99, 6.15, txt_texto3, 22, RGBColor(0, 48, 87))
    agregar_linea(slide, 42.80)
    agregar_texto_avanzado(slide, 1.81, 43.36, 33.99, 12.59, txt_texto4, 22, RGBColor(0, 48, 87))
    agregar_texto_avanzado(slide, 1.87, 50.67, 34.48, 6.96, txt_texto5, 22, RGBColor(0, 48, 87))

    carpeta_salida = "OBERON NEWS/"
    os.makedirs(carpeta_salida, exist_ok=True)

    nombre, extension = os.path.splitext(archivo_salida)
    ruta_final = os.path.join(carpeta_salida, archivo_salida)
    contador = 1
    while os.path.exists(ruta_final):
        ruta_final = os.path.join(carpeta_salida, f"{nombre}_{contador}{extension}")
        contador += 1

    modelo.save(ruta_final)

    messagebox.showinfo("Éxito", f"Boletín generado correctamente.")
    global ultimo_archivo
    ultimo_archivo = ruta_final

# ============================================================
# FUNCIÓN PARA ABRIR ARCHIVO
# ============================================================
def abrir_archivo():
    try:
        if ultimo_archivo and os.path.exists(ultimo_archivo):
            os.system(f'start "" "{ultimo_archivo}"')
            btn_crear_imagen.config(state="normal")
        else:
            messagebox.showwarning("Aviso", "Aún no se ha generado ningún archivo.")
    except Exception as e:
        messagebox.showerror("Error", str(e))

# ============================================================
# FUNCIÓN: CREAR IMAGEN DESDE EL POWERPOINT
# ============================================================
def crear_imagen():
    try:
        if not ultimo_archivo or not os.path.exists(ultimo_archivo):
            messagebox.showwarning("Aviso", "Primero genera un boletín antes de crear la imagen.")
            return

        ruta_pptx = os.path.abspath(ultimo_archivo)

        # Verificar si el archivo está abierto
        try:
            with open(ruta_pptx, "r+b"):
                pass
        except PermissionError:
            messagebox.showwarning(
                "Archivo en uso",
                "⚠️ El archivo PowerPoint está actualmente abierto.\n\n"
                "Por favor ciérralo antes de crear la imagen."
            )
            return

        # Crear carpeta absoluta para imágenes
        carpeta_imagenes = os.path.abspath(os.path.join("OBERON NEWS", "Imagenes"))
        os.makedirs(carpeta_imagenes, exist_ok=True)

        # Generar nombre limpio
        nombre_base = os.path.splitext(os.path.basename(ruta_pptx))[0]
        nombre_limpio = "".join(c for c in nombre_base if c.isalnum() or c in ("_", "-"))
        ruta_salida_jpg = os.path.join(carpeta_imagenes, f"{nombre_limpio}.jpg")

        # Convertir a formato compatible con COM (doble barra)
        ruta_salida_jpg = ruta_salida_jpg.replace("/", "\\")
        ruta_pptx = ruta_pptx.replace("/", "\\")

        # Exportar diapositiva
        powerpoint = win32com.client.Dispatch("PowerPoint.Application")
        powerpoint.Visible = 1
        presentation = powerpoint.Presentations.Open(ruta_pptx)
        presentation.Slides(1).Export(ruta_salida_jpg, "JPG")
        presentation.Close()
        powerpoint.Quit()

        messagebox.showinfo("Éxito", f"Imagen creada correctamente en:\n\n{ruta_salida_jpg}")
        os.startfile(carpeta_imagenes)

    except Exception as e:
        try:
            presentation.Close()
            powerpoint.Quit()
        except:
            pass
        messagebox.showerror("Error", str(e))


# ============================================================
# INTERFAZ RESPONSIVE
# ============================================================
ventana = tk.Tk()
ventana.title("Oberon News – Generador v1.4")
ventana.geometry("1200x680")
ventana.config(bg="#003057")
ventana.resizable(True, True)
ventana._last_text_widget = None



# Sistema de redimensionamiento dinámico
ventana.columnconfigure(0, weight=1)
ventana.rowconfigure(0, weight=1)

# Contenedor principal (divide en columnas izquierda / derecha)
frame_contenedor = tk.Frame(ventana, bg="#003057")
frame_contenedor.grid(row=0, column=0, sticky="nsew")
frame_contenedor.columnconfigure(0, weight=3)  # izquierda
frame_contenedor.columnconfigure(1, weight=2)  # derecha
frame_contenedor.rowconfigure(0, weight=1)

# Frames internos que contendrán tus widgets
frame_izq = tk.Frame(frame_contenedor, bg="#003057")
frame_izq.grid(row=0, column=0, sticky="nsew", padx=20, pady=20)

frame_der = tk.Frame(frame_contenedor, bg="#003057")
frame_der.grid(row=0, column=1, sticky="nsew", padx=20, pady=20)


try:
    logo_img = Image.open(LOGO_EMPRESA)
    logo_img = logo_img.resize((200, 70))
    logo_tk = ImageTk.PhotoImage(logo_img)
    lbl_logo = tk.Label(ventana, image=logo_tk, bg="#003057")
    lbl_logo.place(x=320, y=20)
except:
    pass

titulo_derecha = tk.Label(
    ventana,
    text="¡Oberon News!",
    font=("Century Gothic", 22, "bold"),
    fg="white",
    bg="#003057"
)
titulo_derecha.place(x=900, y=40)

sep = tk.Frame(ventana, bg="white", width=2, height=550)
sep.place(x=820, y=90)

def registrar_texto_activo(event):
    ventana._last_text_widget = event.widget

tk.Label(ventana, text="Fecha:", font=("Century Gothic", 11), fg="white", bg="#003057").place(x=50, y=110)
entry_fecha = tk.Entry(ventana, width=40, font=("Century Gothic", 11))
entry_fecha.insert(0, "21 al 23 de Junio 2025")
entry_fecha.place(x=200, y=110)

tk.Label(ventana, text="Título del boletín:", font=("Century Gothic", 11), fg="white", bg="#003057").place(x=50, y=150)
entry_titulo = tk.Entry(ventana, width=70, font=("Century Gothic", 11))
entry_titulo.insert(0, "AFECTACIONES VIALES Y PLAN DE MOVILIDAD FESTIVO- BOGOTÁ D.C")
entry_titulo.place(x=200, y=150)

def crear_texto(x, y, h, contenido):
    caja = scrolledtext.ScrolledText(ventana, width=80, height=h, font=("Century Gothic", 10))
    caja.insert(tk.END, contenido)
    caja.place(x=x, y=y)
    caja.bind("<FocusIn>", registrar_texto_activo)
    return caja

tk.Label(ventana, text="Texto principal:", font=("Century Gothic", 11), fg="white", bg="#003057").place(x=50, y=190)
txt_texto1 = crear_texto(200, 190, 4,
"Con ocasión del Festival Rock al Parque 2025, que se celebrará los días 21, 22 y 23 de junio en el Parque Simón Bolívar, "
"la Secretaría Distrital de Movilidad anunció una serie de cierres viales y desvíos en sectores clave de la ciudad. "
"Estas medidas, sumadas al incremento en la movilidad por el puente festivo de Corpus Christi y otros eventos programados, "
"como la Marcha del Sur LGTBIQ+, requerirán de una planificación anticipada por parte de la ciudadanía para evitar "
"contratiempos y facilitar el desarrollo seguro y ordenado de las actividades culturales y recreativas en Bogotá.")

tk.Label(ventana, text="Texto 2:", font=("Century Gothic", 11), fg="white", bg="#003057").place(x=50, y=300)
txt_texto2 = crear_texto(200, 300, 4,
"EVENTOS PROGRAMADOS\n\nFestival Rock al Parque 2025\n\n"
"Con motivo del Festival Rock al Parque 2025, que se llevará a cabo el sábado 21, domingo 22 y lunes 23 de junio "
"en el Parque Simón Bolívar, la Secretaría Distrital de Movilidad autorizó cierres y desvíos viales sobre la Av. Calle 63 "
"entre las carreras 60 y 68 en ambos sentidos, entre las 10:30 a. m. y las 11:30 p. m. Asimismo, se verá restringido el uso "
"de la ciclorruta en el mismo tramo y horario.")

tk.Label(ventana, text="Texto 3:", font=("Century Gothic", 11), fg="white", bg="#003057").place(x=50, y=410)
txt_texto3 = crear_texto(200, 410, 3,
"Recomendaciones\n\nSe recomienda a los asistentes al festival y a la ciudadanía en general evitar el uso de vehículo particular y preferir "
"el transporte público (SITP, TransMilenio), programar viajes con anticipación y seguir las indicaciones de las autoridades "
"y de la organización. También se sugiere el uso de transporte no motorizado como la bicicleta.")

tk.Label(ventana, text="Texto 4:", font=("Century Gothic", 11), fg="white", bg="#003057").place(x=50, y=490)
txt_texto4 = crear_texto(200, 490, 3,
"Puente Festivo de Corpus Christi\n\nSe prevé alta movilidad en Bogotá y Cundinamarca durante el puente festivo del 20 al 23 de junio. "
"Se espera la salida de más de 940.000 vehículos desde la capital y el paso de 1,8 millones por los peajes del departamento. "
"El lunes 23 de junio se activará el reversible en la carrera 7 entre calles 245 y 183, y se implementará pico y placa regional: "
"placas pares de 12:00 m. a 4:00 p.m. y impares de 4:00 p.m. a 8:00 p.m.")

tk.Label(ventana, text="Texto 5:", font=("Century Gothic", 11), fg="white", bg="#003057").place(x=50, y=570)
txt_texto5 = crear_texto(200, 570, 3,
"Marcha del Sur LGTBIQ+\n\nEl domingo 22 de junio a partir de las 7:00 a. m., se realizará la XVII Marcha del Sur LGTBIQ+, una manifestación respaldada "
"por la Alcaldía Mayor de Bogotá. El evento contará con un recorrido que iniciará en la Plazoleta Fundacional de Bosa "
"(carrera 86 con calle 1.° de mayo) y culminará en la Alcaldía Local de Kennedy. La actividad es convocada por diversos "
"colectivos LGTBIQ+.")

# ==============================
# COLUMNA DERECHA ALINEADA
# ==============================
def seleccionar_imagen():
    archivo = filedialog.askopenfilename(title="Seleccionar imagen evidencia",
                                         filetypes=[("Imágenes", "*.png *.jpg")])
    if archivo:
        entry_imagen.delete(0, tk.END)
        entry_imagen.insert(0, archivo)

col_x = 870
y_base = 180
espacio = 40

tk.Label(ventana, text="Imagen evidencia (PNG/JPG):", font=("Century Gothic", 11),
         fg="white", bg="#003057").place(x=col_x, y=y_base)
entry_imagen = tk.Entry(ventana, width=35, font=("Century Gothic", 10), justify="center")
entry_imagen.insert(0, "EVIDENCIAS/imagen2.png")
entry_imagen.place(x=col_x, y=y_base+30)
ttk.Button(ventana, text="Examinar...", command=seleccionar_imagen, width=20).place(x=col_x+40, y=y_base+60)

tk.Label(ventana, text="Nombre del archivo (.pptx):", font=("Century Gothic", 11),
         fg="white", bg="#003057").place(x=col_x, y=y_base+120)
entry_nombre = tk.Entry(ventana, width=35, font=("Century Gothic", 10), justify="center")
nombre_defecto = f"ON {datetime.now().strftime('%d-%m-%Y')}.pptx"
entry_nombre.insert(0, nombre_defecto)
entry_nombre.place(x=col_x, y=y_base+150)

def _aplicar_formato(tag, **config):
    text_widget = ventana._last_text_widget if isinstance(getattr(ventana, "_last_text_widget", None), tk.Text) else None
    if not text_widget:
        for tw in (txt_texto1, txt_texto2, txt_texto3, txt_texto4, txt_texto5):
            try:
                _ = tw.index("sel.first")
                text_widget = tw
                break
            except tk.TclError:
                continue
    if not text_widget:
        messagebox.showinfo("Aviso", "Selecciona el texto dentro de una caja antes de aplicar formato.")
        return
    try:
        start, end = text_widget.index("sel.first"), text_widget.index("sel.last")
        text_widget.tag_configure(tag, **config)
        text_widget.tag_add(tag, start, end)
    except tk.TclError:
        messagebox.showinfo("Aviso", "Selecciona primero el texto que deseas modificar.")

def aplicar_negrilla(): _aplicar_formato("bold", font=("Century Gothic", 10, "bold"))
def aplicar_rojo_oberon(): _aplicar_formato("rojo", foreground="#D64123")
def aplicar_azul_oberon(): _aplicar_formato("azul", foreground="#003057")

# ==============================
# ESTILO DE BOTONES OBERON
# ==============================
style = ttk.Style()
style.theme_use("clam")  # tema que permite modificar colores

# Estilo general para todos los botones
style.configure(
    "TButton",
    font=("Century Gothic", 10, "bold"),
    foreground="white",
    background="#00A3AF",
    borderwidth=2,
    focusthickness=3,
    focuscolor="none"
)

# Colores dinámicos (hover, deshabilitado)
style.map(
    "TButton",
    background=[("active", "#008C97"), ("disabled", "#00727B")],
    foreground=[("active", "white"), ("disabled", "#CCCCCC")]
)

# Estilo especial solo para el botón “Generar Boletín”
style.configure(
    "Generar.TButton",
    background="#D64123",
    foreground="white"
)
style.map(
    "Generar.TButton",
    background=[("active", "#B8321B"), ("disabled", "#A22C18")],
    foreground=[("active", "white"), ("disabled", "#EEEEEE")]
)


y_boton = y_base + 200

# Botones turquesa (por defecto usan el estilo "TButton")
ttk.Button(ventana, text="✏️ Negrilla", command=aplicar_negrilla, width=20).place(x=col_x+40, y=y_boton)
ttk.Button(ventana, text="🔴 Rojo Oberon", command=aplicar_rojo_oberon, width=20).place(x=col_x+40, y=y_boton+40)
ttk.Button(ventana, text="🔵 Azul Oberon", command=aplicar_azul_oberon, width=20).place(x=col_x+40, y=y_boton+80)
ttk.Button(ventana, text="📂 Abrir archivo", command=abrir_archivo, width=20).place(x=col_x+40, y=y_boton+160)
btn_crear_imagen = ttk.Button(ventana, text="🖼️ Crear Imagen", command=crear_imagen, width=20, state="disabled")
btn_crear_imagen.place(x=col_x+40, y=y_boton+200)
ttk.Button(ventana, text="❌ Salir", command=ventana.destroy, width=20).place(x=col_x+40, y=y_boton+240)

# Botón rojo Oberon (usa estilo especial)
ttk.Button(
    ventana,
    text="🧾 Generar Boletín",
    style="Generar.TButton",   
    command=generar_boletin,
    width=20
).place(x=col_x+40, y=y_boton+120)


ultimo_archivo = None
ventana.mainloop()
