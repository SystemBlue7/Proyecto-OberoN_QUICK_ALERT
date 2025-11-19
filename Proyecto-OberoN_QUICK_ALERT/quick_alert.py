import os
import locale
from datetime import datetime
import tkinter as tk
from tkinter import ttk, filedialog, scrolledtext, messagebox
from pptx import Presentation
from pptx.util import Cm, Pt
from PIL import Image, ImageTk
import win32com.client
import pyodbc

# FUNCIÓN PARA RUTAS CORRECTAS EN .PY Y .EXE
def resource_path(relative_path):
    import sys, os
    if hasattr(sys, "_MEIPASS"):
        return os.path.join(sys._MEIPASS, relative_path)
    return os.path.join(os.path.abspath("."), relative_path)

# ---------------------------------------------
# MAPA CIUDAD → MUNICIPIOS
# ---------------------------------------------
MUNICIPIOS_POR_CIUDAD = {
    "Bogotá": [
        "Bogotá", "Soacha", "Chía", "Zipaquirá", "Cajicá", "Funza", "Mosquera",
        "Madrid", "Facatativá", "La Calera", "Sopó", "Cota", "Tenjo",
        "Guasca", "Guatavita", "Tabio", "Subachoque"
    ],
    "Medellín": [
        "Medellín", "Envigado", "Itagüí", "Bello", "Sabaneta", "La Estrella",
        "Caldas", "Copacabana", "Girardota", "Barbosa"
    ],
    "Cali": [
        "Cali", "Palmira", "Yumbo", "Jamundí", "Candelaria", "Florida",
        "Pradera"
    ],
    "Barranquilla": [
        "Barranquilla", "Soledad", "Malambo", "Galapa", "Puerto Colombia"
    ],
    "Cartagena": [
        "Cartagena", "Turbaco", "Arjona", "Santa Rosa", "Turbana", "Clemencia"
    ],
    "Cúcuta": [
        "Cúcuta", "Villa del Rosario", "Los Patios", "El Zulia", "San Cayetano"
    ],
    "Bucaramanga": [
        "Bucaramanga", "Floridablanca", "Girón", "Piedecuesta", "Lebrija"
    ],
    "Pereira": [
        "Pereira", "Dosquebradas", "La Virginia", "Santa Rosa de Cabal"
    ],
    "Santa Marta": [
        "Santa Marta", "Ciénaga", "Zona Bananera"
    ],
    "Ibagué": [
        "Ibagué", "Espinal", "Cajamarca", "Alvarado", "Piedras"
    ],
    "Villavicencio": [
        "Villavicencio", "Acacías", "Restrepo", "Cumaral", "Guamal"
    ],
    "Tunja": [
        "Tunja", "Sogamoso", "Duitama", "Paipa"
    ],
    "Valledupar": [
        "Valledupar", "La Paz", "Codazzi", "San Diego", "Manaure"
    ],
    "Florencia": [
        "Florencia", "La Montañita", "Morelia", "Milán"
    ],
    "Quibdó": [
        "Quibdó", "Lloró", "Yuto", "Bojayá"
    ],
    "Neiva": [
        "Neiva", "Rivera", "Campoalegre", "Tello", "Aipe"
    ],
    "Manizales": [
        "Manizales", "Villamaría", "Chinchiná", "Neira"
    ],
    "Pasto": [
        "Pasto", "Tangua", "Yacuanquer", "Chachagüí"
    ],
    "Montería": [
        "Montería", "Cereté", "San Pelayo", "Ciénaga de Oro"
    ],
    "Sincelejo": [
        "Sincelejo", "Corozal", "Morroa", "Los Palmitos"
    ],
    "Popayán": [
        "Popayán", "Timbío", "Cajibío", "El Tambo"
    ]
}

CIUDADES_COLOMBIA = list(MUNICIPIOS_POR_CIUDAD.keys())



# ----------------------------------------------------------
# CONFIGURACIÓN DE IDIOMA Y FECHA
# ----------------------------------------------------------
for loc in ("es_ES.UTF-8", "es_CO.UTF-8", "es_ES.utf8", "Spanish_Colombia", "Spanish_Spain"):
    try:
        locale.setlocale(locale.LC_TIME, loc)
        break
    except Exception:
        continue

fecha_formateada = datetime.now().strftime("%d de %B del %Y")

# ----------------------------------------------------------
# FUNCIONES DE TEXTO
# ----------------------------------------------------------
def tk_index_to_flat(widget, index):
    linea, columna = map(int, str(index).split("."))
    texto = widget.get("1.0", "end-1c")
    lineas = texto.split("\n")
    prev = sum(len(lineas[i]) + 1 for i in range(linea - 1))
    return prev + columna

def get_bold_ranges(widget):
    if "bold" not in widget.tag_names():
        return []
    ranges = widget.tag_ranges("bold")
    if not ranges:
        return []
    planos = [tk_index_to_flat(widget, widget.index(r)) for r in ranges]
    return list(zip(planos[::2], planos[1::2]))

def _prep_textframe(tf):
    tf.margin_top = tf.margin_bottom = tf.margin_left = tf.margin_right = 0
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = ""
    p.space_before = p.space_after = 0
    p.line_spacing = 1.0
    return p

def add_rich_text_box(slide, x, y, w, h, texto, size_pt=23, bold_ranges=None):
    if bold_ranges is None:
        bold_ranges = []
    bold_ranges = sorted(bold_ranges, key=lambda t: t[0])
    tb = slide.shapes.add_textbox(Cm(x), Cm(y), Cm(w), Cm(h))
    p = _prep_textframe(tb.text_frame)
    last = 0
    n = len(texto)
    for (ini, fin) in bold_ranges:
        ini = max(0, min(ini, n))
        fin = max(0, min(fin, n))
        if ini > last:
            r = p.add_run()
            r.text = texto[last:ini]
            r.font.name = "Century Gothic"
            r.font.size = Pt(size_pt)
            r.font.bold = False
        if fin > ini:
            rb = p.add_run()
            rb.text = texto[ini:fin]
            rb.font.name = "Century Gothic"
            rb.font.size = Pt(size_pt)
            rb.font.bold = True
        last = max(last, fin)
    if last < n:
        r = p.add_run()
        r.text = texto[last:n]
        r.font.name = "Century Gothic"
        r.font.size = Pt(size_pt)
        r.font.bold = False
    return tb.text_frame

def add_simple_text(slide, x, y, w, h, texto, size_pt, negrita=False):
    tb = slide.shapes.add_textbox(Cm(x), Cm(y), Cm(w), Cm(h))
    p = _prep_textframe(tb.text_frame)
    run = p.add_run()
    run.text = texto
    run.font.name = "Century Gothic"
    run.font.size = Pt(size_pt)
    run.font.bold = negrita
    return tb.text_frame

# ----------------------------------------------------------
# CONVERSIÓN PPTX → PNG
# ----------------------------------------------------------
def convertir_pptx_a_png(ruta_pptx):
    try:
        if os.path.exists(ruta_pptx):
            try:
                os.rename(ruta_pptx, ruta_pptx)
            except PermissionError:
                messagebox.showwarning("Archivo en uso", "Cierra el archivo PowerPoint antes de convertir a PNG.")
                return

        powerpoint = win32com.client.Dispatch("PowerPoint.Application")
        powerpoint.Visible = True
        powerpoint.WindowState = 2

        pres = powerpoint.Presentations.Open(ruta_pptx, WithWindow=False)

        carpeta_salida = os.path.dirname(ruta_pptx)
        nombre_base = os.path.splitext(os.path.basename(ruta_pptx))[0]
        carpeta_png = os.path.join(carpeta_salida, "IMAGEN PNG")
        os.makedirs(carpeta_png, exist_ok=True)

        existentes_antes = set(f for f in os.listdir(carpeta_png) if f.lower().endswith(".png"))

        pres.Export(carpeta_png, "PNG")
        pres.Close()
        powerpoint.Quit()

        import time, shutil

        tiempo_maximo = 8
        inicio = time.time()
        nuevos_pngs = []

        while time.time() - inicio < tiempo_maximo:
            despues = [f for f in os.listdir(carpeta_png) if f.lower().endswith(".png")]
            nuevos_pngs = [f for f in despues if f not in existentes_antes]
            if nuevos_pngs:
                break
            time.sleep(0.25)

        if not nuevos_pngs:
            messagebox.showwarning("Aviso", "No se generó la imagen PNG. Intenta nuevamente.")
            return

        nuevos_pngs.sort(key=lambda x: x.lower())
        archivo_png = os.path.join(carpeta_png, nuevos_pngs[0])

        destino = os.path.join(carpeta_png, f"{nombre_base}.png")
        contador = 2
        while os.path.exists(destino):
            destino = os.path.join(carpeta_png, f"{nombre_base} ({contador}).png")
            contador += 1

        shutil.move(archivo_png, destino)

        for f in nuevos_pngs[1:]:
            ruta_f = os.path.join(carpeta_png, f)
            try:
                os.remove(ruta_f)
            except:
                pass

        os.startfile(destino)

    except Exception as e:
        messagebox.showerror("Error", f"No se pudo convertir a PNG:\n{e}")

# ----------------------------------------------------------
# CLASE PRINCIPAL
# ----------------------------------------------------------
class AppBoletin(tk.Tk):
    def __init__(self):
        super().__init__()
        self.title("Boletín Oberon (v10.0)")
        self.geometry("900x768")

        # Ícono personalizado compatible con EXE
        # Ícono personalizado compatible con ejecución normal y .exe
        try:
            import sys, os

            if getattr(sys, 'frozen', False):
                # Si corre como EXE
                icon_path = os.path.join(sys._MEIPASS, "O-PEQUEÑA.ico")
            else:
                # Si corre como script .py
                icon_path = os.path.join(os.path.dirname(__file__), "O-PEQUEÑA.ico")

            self.iconbitmap(icon_path)

        except Exception as e:
            print("ERROR ICONO:", e)
            pass


        self._vars()
        self._create_styles()
        self._apply_theme()
        self.ultimo_archivo = None
        self.icon_previews = {}
        self._all_frames = []
        self._last_text_widget = None
        self._make_responsive()
        self._ui()


    # ----------------- VARIABLES -----------------
    def _vars(self):
        self.font_ui = ("Century Gothic", 10)
        self.font_title = ("Century Gothic", 18, "bold")

        self.var_icono_ubicacion = tk.StringVar(value="")
        self.var_icono_titulo = tk.StringVar(value="")
        self.var_icono_boletin = tk.StringVar(value="")
        self.var_icono_recomendaciones = tk.StringVar(value="")

        self.var_ubicacion = tk.StringVar(value="CALLE 45")
        self.var_titulo = tk.StringVar(value="BLOQUEO POR PLANTÓN DURANTE AUDIENCIA JUDICIAL DEL EXPRESIDENTE")
        self.var_boletin = tk.StringVar(value="Boletín No.")
        self.var_fuente = tk.StringVar(value="Fuente: W Radio")

        self.var_impacto = tk.StringVar(value="Impacto del evento en la movilidad y la seguridad vial.")
        
        self.var_ciudad = tk.StringVar(value="Bogotá")
        self.var_municipio = tk.StringVar(value="Bogotá")

        

        self.var_cabezote = tk.StringVar(value="CABEZOTES/MOVILIDAD.JPG")
        self.var_evidencia = tk.StringVar(value="EVIDENCIAS/imagen_1.png")
        self.var_nombre_archivo = tk.StringVar(
            value=f"OQ {datetime.now().strftime('%d-%m-%Y')}.pptx"
        )

    def _clear_entry_on_click(self, event):
        widget = event.widget
        if not isinstance(widget, ttk.Entry):
            return

        texto_actual = widget.get().strip()

        # Si está vacío, restaurar el placeholder
        if texto_actual == "":
            widget.delete(0, tk.END)
            widget.insert(0, "Boletín No.")
            return

        # Si tiene el placeholder, no borrar
        if texto_actual == "Boletín No.":
            return

        # Si tiene texto del usuario, borrarlo
        widget.delete(0, tk.END)



    def _clear_text_on_click(self, event):
        widget = event.widget
        if isinstance(widget, scrolledtext.ScrolledText):
            widget.delete("1.0", tk.END)
            
    def _actualizar_municipios(self, *args):
        ciudad = self.var_ciudad.get()
        municipios = MUNICIPIOS_POR_CIUDAD.get(ciudad, [])

        self.combo_municipio["values"] = municipios

        if municipios:
            self.var_municipio.set(municipios[0])
        else:
            self.var_municipio.set("")


    # ----------------- ESTILOS -----------------
    def _create_styles(self):
        self.style = ttk.Style(self)
        try:
            self.style.theme_use("clam")
        except Exception:
            pass

        self.palette_dark = {
            "bg": "#012B47",
            "fg": "#FFFFFF",
            "entry_bg": "#0F3E5E",
            "entry_fg": "#FFFFFF",
            "entry_insert": "#FFFFFF",
            "button_primary": "#007D7D",
            "button_accent": "#D64123",
            "button_text": "#FFFFFF",
            "placeholder": "#A9A9A9"
        }

        for sty in ("TFrame", "TLabel", "TEntry"):
            self.style.configure(sty, font=self.font_ui)

        self.style.configure("Primary.TButton", font=("Century Gothic", 9, "bold"), padding=(4, 2))
        self.style.map("Primary.TButton", background=[("active", "!disabled", "#0A9F9F")])
        self.style.configure("Accent.TButton", font=("Century Gothic", 9, "bold"), padding=(4, 2))
        self.style.map("Accent.TButton", background=[("active", "!disabled", "#E35A3A")])

        # =======================
        # ESTILO REAL PARA COMBOBOX OSCURO (CIUDAD Y MUNICIPIO)
        # =======================

        # Importa la flecha del tema clam SIN cambiar todo el sistema
        self.style.element_create("CustomCombobox.downarrow", "from", "clam")

        self.style.layout("Lugar.TCombobox", [
            ('CustomCombobox.downarrow', {'side': 'right', 'sticky': ''}),
            ('Combobox.padding',   {'children': [
                ('Combobox.textarea', {'sticky': 'nswe'})
            ]})
        ])

        self.style.configure(
            "Lugar.TCombobox",
            fieldbackground="#003057",
            background="#003057",
            foreground="#FFFFFF",
            arrowcolor="#FFFFFF"
        )

        self.option_add('*TCombobox*Listbox*Background', '#003057')
        self.option_add('*TCombobox*Listbox*Foreground', '#FFFFFF')



    def _apply_theme(self):
        pal = self.palette_dark
        self.configure(bg=pal["bg"])
        self.style.configure("TFrame", background=pal["bg"])
        self.style.configure("TLabel", background=pal["bg"], foreground=pal["fg"])
        self.style.configure("TEntry",
                             fieldbackground=pal["entry_bg"],
                             foreground=pal["entry_fg"],
                             insertcolor=pal["entry_insert"])
        self.style.configure("Primary.TButton", background=pal["button_primary"], foreground=pal["button_text"])
        self.style.configure("Accent.TButton", background=pal["button_accent"], foreground=pal["button_text"])

    # ----------------- AUXILIARES -----------------
    def _fila_entry(self, parent, text, var):
        f = ttk.Frame(parent)
        f.pack(fill=tk.X, pady=4)

        ttk.Label(f, text=text).pack(side=tk.LEFT)
        entry = ttk.Entry(f, textvariable=var, foreground="#A9A9A9")
        entry.pack(side=tk.LEFT, fill=tk.X, expand=True, padx=6)

        entry.bind("<Button-1>", self._clear_entry_on_click)

        return f

    def _selector_vertical(self, parent, text, var, row):
        f = ttk.Frame(parent)
        f.grid(row=row, column=0, pady=3)
        ttk.Label(f, text=text).pack(anchor="w")
        ttk.Entry(f, textvariable=var, state="readonly", foreground="#A9A9A9").pack(fill=tk.X, pady=2)
        ttk.Button(f, text="Examinar...", style="Primary.TButton",
                   command=lambda: self._examinar(var)).pack(pady=2)
        return f
    
    def _validar_requisitos(self):
        requisitos = [
            self.var_icono_ubicacion.get().strip(),
            self.var_icono_titulo.get().strip(),
            self.var_icono_boletin.get().strip(),
            self.var_icono_recomendaciones.get().strip(),
            self.var_cabezote.get().strip(),
            self.var_evidencia.get().strip()
        ]

        # Verificar que todos los archivos existan realmente
        completos = all(os.path.isfile(r) for r in requisitos)

        if completos:
            self.btn_generar.state(["!disabled"])
        else:
            self.btn_generar.state(["disabled"])


    # ----------------- INTERFAZ -----------------
    def _ui(self):
        cont = ttk.Frame(self)
        cont.pack(fill=tk.BOTH, expand=True, padx=12, pady=12)

        cont.columnconfigure(0, weight=1)
        cont.columnconfigure(1, weight=1)
        cont.rowconfigure(0, weight=1)

        left = ttk.Frame(cont)
        left.grid(row=0, column=0, sticky="nsew", padx=(0, 10))

        right = ttk.Frame(cont)
        right.grid(row=0, column=1, sticky="nsew")

        self._all_frames.extend([cont, left, right])

        try:
            logo_path = os.path.join(os.path.dirname(__file__), "Logo.png")
            logo_img = Image.open(logo_path)
            logo_img.thumbnail((200, 200))
            self.tk_logo = ImageTk.PhotoImage(logo_img)
            ttk.Label(left, image=self.tk_logo, background="#012B47").pack(anchor="center", pady=(0, 0))
        except Exception:
            pass
        
        # ----- Ciudad y Municipio en una sola línea -----
        frame_lugares = ttk.Frame(left)
        frame_lugares.pack(fill=tk.X, pady=(6, 7))

        # Ciudad
        ttk.Label(frame_lugares, text="Ciudad:").grid(row=0, column=0, sticky="w", padx=(110, 5))

        self.combo_ciudad = ttk.Combobox(
            frame_lugares,
            textvariable=self.var_ciudad,
            values=list(MUNICIPIOS_POR_CIUDAD.keys()),
            state="readonly",
            width=22,
            style="Lugar.TCombobox"
        )


        self.combo_ciudad.grid(row=0, column=1, sticky="w", pady=2)
        self.combo_ciudad.bind("<<ComboboxSelected>>", self._actualizar_municipios)


        # Municipio
        ttk.Label(frame_lugares, text="Municipio:").grid(row=0, column=2, sticky="w", padx=(50, 5))

        self.combo_municipio = ttk.Combobox(
            frame_lugares,
            textvariable=self.var_municipio,
            values=MUNICIPIOS_POR_CIUDAD.get(self.var_ciudad.get(), []),
            state="readonly",
            width=22,
            style="Lugar.TCombobox"
        )

        self.combo_municipio.grid(row=0, column=3, sticky="w", pady=2)

        self.f_ubic = self._fila_entry(left, "Ubicación:", self.var_ubicacion)
        ttk.Button(left, text="Seleccionar ícono (Ubicación)", style="Primary.TButton",
                   command=self._seleccionar_icono_ubicacion).pack(pady=4)

        self.f_tit = self._fila_entry(left, "Título:", self.var_titulo)
        ttk.Button(left, text="Seleccionar ícono (Título)", style="Primary.TButton",
                   command=self._seleccionar_icono_titulo).pack(pady=4)

        self.f_bol = self._fila_entry(left, "Número de Boletín:", self.var_boletin)
        ttk.Button(left, text="Seleccionar ícono (Boletín)", style="Primary.TButton",
                   command=self._seleccionar_icono_boletin).pack(pady=4)

        ttk.Label(left, text="Texto principal:").pack(anchor="w", pady=(8, 2))
        self.txt_principal = scrolledtext.ScrolledText(left, height=6, wrap=tk.WORD, fg="#A9A9A9")
        self.txt_principal.bind("<Button-1>", self._clear_text_on_click)
        self.txt_principal.pack(fill=tk.X)
        self.txt_principal.insert(tk.END,
            "Este 19 de octubre de 2025, en la tarde de este lunes se registra un bloqueo en la Carrera 7 con Calle 45, "
            "en sentido occidente-oriente, generado por un plantón convocado por el colectivo “Puro Veneno” con motivo de la "
            "audiencia judicial del expresidente Álvaro Uribe Vélez. Este 28 de julio de 2025, en la tarde de este lunes se "
            "registra un bloqueo en la Carrera 7 con Calle 45, en sentido occidente-oriente, generado por un plantón convocado "
            "por el colectivo “Puro Veneno” con motivo de la audiencia judicial del expresidente Álvaro Uribe Vélez."
        )

        ttk.Label(left, text="Texto secundario:").pack(anchor="w", pady=(8, 2))
        self.txt_secundario = scrolledtext.ScrolledText(left, height=6, wrap=tk.WORD, fg="#A9A9A9")
        self.txt_secundario.bind("<Button-1>", self._clear_text_on_click)
        self.txt_secundario.pack(fill=tk.X)
        self.txt_secundario.insert(tk.END,
            "La manifestación ha generado afectación de la calzada mixta, obligando a las autoridades "
            "a cerrar la Calle 45 entre Carreras 13 y 7. Unidades de tránsito y policía hacen presencia "
            "en la zona para controlar la situación y garantizar la seguridad vial."
        )
        
        # Campos Fuente e Impacto
        self._fila_entry(left, "Fuente:", self.var_fuente)
        self._fila_entry(left, "Impacto:", self.var_impacto)


        
        self.frame_recom = ttk.Frame(left)
        self.frame_recom.pack(anchor="w", pady=(0, 2))

        self.btn_recom = ttk.Button(self.frame_recom, text="Seleccionar ícono (Recomendaciones)",
                                    style="Primary.TButton", command=self._seleccionar_icono_recomendaciones)
        self.btn_recom.pack(side=tk.LEFT, anchor="w")

        
       
        self.lbl_recom_icon = ttk.Label(self.frame_recom)
        self.lbl_recom_icon.pack(side=tk.RIGHT, padx=6, pady=4)

        self.txt_recomendaciones = scrolledtext.ScrolledText(left, height=10, wrap=tk.WORD, fg="#A9A9A9")
        self.txt_recomendaciones.bind("<Button-1>", self._clear_text_on_click)
        self.txt_recomendaciones.pack(fill=tk.BOTH, expand=True)
        self.txt_recomendaciones.insert(tk.END,
            "✅Ante los bloqueos reportados en diferentes puntos de la ciudad, como la Carrera 7 con Calle 45 "
            "y la Av. Guayacanes con Calle 38 Sur, las autoridades emiten las siguientes recomendaciones:\n\n"
            "✅Evite transitar por zonas de manifestación: si puede, reprograme sus desplazamientos o utilice rutas alternas "
            "como la Avenida Caracas, Carrera 30 (NQS) o Avenida Circunvalar.\n\n"
            "✅Proteja a menores y adultos mayores: si se moviliza con personas vulnerables, evite transitar cerca de "
            "aglomeraciones por seguridad.\n\n"
            "✅Use transporte alternativo: de ser posible, opte por el uso de bicicleta, caminar si la distancia lo permite, "
            "o planifique rutas con aplicaciones de movilidad en tiempo real."
        )

        for tw in (self.txt_principal, self.txt_secundario, self.txt_recomendaciones):
            tw.bind("<FocusIn>", self._remember_text_widget)
            tw.bind("<Button-1>", self._remember_text_widget)
            tw.bind("<Key>", self._remember_text_widget)

        for i in range(12):
            right.grid_rowconfigure(i, weight=1)
        right.grid_columnconfigure(0, weight=1)

        self.lbl_right_title = ttk.Label(right, text="Quick Alert", font=self.font_title,
                                         anchor="center", foreground="#FFFFFF")
        self.lbl_right_title.grid(row=0, column=0, pady=(5, 10))

        self._selector_vertical(right, "Cabezote (JPG):", self.var_cabezote, 1)
        self._selector_vertical(right, "Imagen evidencia (PNG/JPG):", self.var_evidencia, 2)

        ttk.Label(right, text="Nombre del archivo (.pptx):").grid(row=3, column=0)
        ttk.Entry(right, textvariable=self.var_nombre_archivo, foreground="#A9A9A9").grid(row=4, column=0, padx=10, pady=3)

        ttk.Button(right, text="Aplicar Negrilla", style="Primary.TButton",
                   command=self._aplicar_negrilla).grid(row=5, column=0, pady=3)

        self.btn_generar = ttk.Button(right, text="Generar Boletín PPTX",
                              style="Accent.TButton", command=self._crear)
        self.btn_generar.grid(row=6, column=0, pady=3)
        self.btn_generar.state(["disabled"])


        ttk.Button(right, text="Abrir archivo generado", style="Primary.TButton",
                   command=self._abrir_archivo).grid(row=7, column=0, pady=3)

        self.btn_png = ttk.Button(right, text="Convertir a PNG", style="Primary.TButton",
                                  command=self._convertir_a_png)
        self.btn_png.grid(row=8, column=0, pady=3)
        self.btn_png.state(["disabled"])

        self.btn_sql = ttk.Button(right, text="Cargar a SQL", style="Primary.TButton",
                          command=self._cargar_sql)
        self.btn_sql.grid(row=9, column=0, pady=3)
        self.btn_sql.state(["disabled"])    


        ttk.Button(right, text="Salir", style="Primary.TButton",
                   command=self.destroy).grid(row=10, column=0, pady=3)

    # ----------------- FUNCIONES RESTANTES -----------------
    def _examinar(self, var):
        ruta = filedialog.askopenfilename(
            title="Seleccionar imagen",
            filetypes=[("Imágenes", "*.png;*.jpg;*.jpeg")]
        )
        if not ruta:
            return

        ext = os.path.splitext(ruta)[1].lower()
        if ext not in [".png", ".jpg", ".jpeg"]:
            messagebox.showwarning(
            "Formato no permitido",
            "Solo se permiten imágenes en formato PNG o JPG.\nPor favor selecciona otro archivo."
        )
            return

        var.set(os.path.abspath(ruta))
        self._validar_requisitos()

    def _make_responsive(self):
        self.rowconfigure(0, weight=1)
        self.columnconfigure(0, weight=1)

    def _seleccionar_icono_generico(self, variable, frame_ref):
        ruta = filedialog.askopenfilename(
            title="Seleccionar ícono",
            initialdir="ICONOGRAFIA",
            filetypes=[("Imágenes", "*.png;*.jpg;*.jpeg;*.gif")]
        )
        if not ruta:
            return
        variable.set(ruta)
        try:
            img = Image.open(ruta)
            img.thumbnail((24, 24))
            tk_img = ImageTk.PhotoImage(img)
            lbl = self.icon_previews.get(frame_ref)
            if lbl is None:
                lbl = ttk.Label(frame_ref)
                lbl.pack(side=tk.RIGHT, padx=4)
                self.icon_previews[frame_ref] = lbl
            lbl.configure(image=tk_img)
            lbl.image = tk_img
        except:
            pass

    def _seleccionar_icono_ubicacion(self):
        self._seleccionar_icono_generico(self.var_icono_ubicacion, self.f_ubic)
        self._validar_requisitos()

    def _seleccionar_icono_titulo(self):
        self._seleccionar_icono_generico(self.var_icono_titulo, self.f_tit)
        self._validar_requisitos()

    def _seleccionar_icono_boletin(self):
        self._seleccionar_icono_generico(self.var_icono_boletin, self.f_bol)
        self._validar_requisitos()


    def _seleccionar_icono_recomendaciones(self):
        ruta = filedialog.askopenfilename(
            title="Seleccionar ícono",
            initialdir="ICONOGRAFIA",
            filetypes=[("Imágenes", "*.png;*.jpg;*.jpeg;*.gif")]
        )
        if ruta:
            self.var_icono_recomendaciones.set(ruta)
            try:
                img = Image.open(ruta)
                img.thumbnail((24, 24))
                tk_img = ImageTk.PhotoImage(img)
                self.lbl_recom_icon.configure(image=tk_img)
                self.lbl_recom_icon.image = tk_img
            except:
                pass

        self._validar_requisitos()


    def _remember_text_widget(self, event):
        self._last_text_widget = event.widget

    def _aplicar_negrilla(self):
        text_widget = None
        if isinstance(self._last_text_widget, tk.Text):
            text_widget = self._last_text_widget
        else:
            widget = self.focus_get()
            if isinstance(widget, tk.Text):
                text_widget = widget

        if text_widget is None:
            for tw in (getattr(self, 'txt_principal', None),
                       getattr(self, 'txt_secundario', None),
                       getattr(self, 'txt_recomendaciones', None)):
                if isinstance(tw, scrolledtext.ScrolledText):
                    try:
                        _ = tw.index("sel.first")
                        text_widget = tw
                        break
                    except tk.TclError:
                        continue

        if text_widget is None:
            messagebox.showinfo("Aviso", "Selecciona texto para aplicar negrilla.")
            return

        try:
            start = text_widget.index("sel.first")
            end = text_widget.index("sel.last")
            text_widget.tag_configure("bold", font=("Century Gothic", 10, "bold"))
            text_widget.tag_add("bold", start, end)
        except tk.TclError:
            messagebox.showinfo("Aviso", "Selecciona primero el texto para aplicar negrilla.")

    def _abrir_archivo(self):
        if self.ultimo_archivo and os.path.exists(self.ultimo_archivo):
            os.startfile(self.ultimo_archivo)
            self.btn_png.state(["!disabled"])
        else:
            messagebox.showwarning("Atención", "No hay archivo generado o no se encuentra.")

    def _convertir_a_png(self):
        if self.ultimo_archivo and os.path.exists(self.ultimo_archivo):
            convertir_pptx_a_png(self.ultimo_archivo)
        try:
            self.btn_sql.state(["!disabled"])
        except:
            pass
        else:
            return
        
    def _cargar_sql(self):
        
        try:
            # =============================================
            # 1. Conexión a Azure SQL
            # =============================================
            conn = pyodbc.connect(
                "DRIVER={ODBC Driver 17 for SQL Server};"
                "SERVER=dwoberon.database.windows.net;"
                "DATABASE=DWH_OBERON_360;"
                "UID=oberon;"
                "PWD=G7p!xQ2v#Lm9zT@w;"
            )

            cursor = conn.cursor()

            # =============================================
            # 2. Captura de datos desde la interfaz
            # =============================================
            fecha = datetime.now().strftime("%Y-%m-%d %H:%M:%S")   # FECHA REAL SQL
            ubicacion = self.var_ubicacion.get()
            titulo = self.var_titulo.get()

            # Ruta ORIGINAL para obtener el nombre correcto (no cambia en EXE)
            ruta_original = self.var_cabezote.get()
            print("DEBUG NOMBRE ARCHIVO:", ruta_original)


            # Nombre sin extensión
            nombre_cabezote = os.path.splitext(os.path.basename(ruta_original))[0]
            print("DEBUG NOMBRE SIN EXT:", nombre_cabezote)

            # === ELIMINAR SOLO LA PALABRA "BLANCO" ===
            nombre_cabezote = nombre_cabezote.replace("BLANCO", "")
            nombre_cabezote = nombre_cabezote.replace("Blanco", "")
            nombre_cabezote = nombre_cabezote.replace("blanco", "")

            # Limpieza de espacios dobles que puedan quedar
            nombre_cabezote = " ".join(nombre_cabezote.split())



            # Quitar la última parte separada por "_"
            partes = nombre_cabezote.split("_")
            if len(partes) > 1:
                nombre_cabezote = "_".join(partes[:-1])

            # Ruta CORRECTA para que el EXE acceda al archivo
            ruta_cabezote = resource_path(ruta_original)

            texto_principal = self.txt_principal.get("1.0", "end-1c")
            impacto = self.var_impacto.get()
            numero_boletin = self.var_boletin.get()
            ciudad = self.var_ciudad.get()
            municipio = self.var_municipio.get()

            # =============================================
            # 3. Inserción de datos
            # =============================================
            cursor.execute("""
                INSERT INTO HISTORICO_QUICK_ALERT (
                    fecha,
                    ubicacion,
                    titulo,
                    nombre_cabezote,
                    texto_principal,
                    impacto,
                    numero_boletin,
                    ciudad,
                    municipio
                )
                VALUES (?, ?, ?, ?, ?, ?, ?, ?, ?)
            """, (
                fecha,
                ubicacion,
                titulo,
                nombre_cabezote,
                texto_principal,
                impacto,
                numero_boletin,
                ciudad,
                municipio
            ))

            conn.commit()
            conn.close()

            messagebox.showinfo("Éxito", "Datos cargados correctamente en SQL.")

        except Exception as e:
            messagebox.showerror("Error SQL", f"Ocurrió un error al cargar datos:\n{e}")


    # ----------------------------------------------------------
    # CREAR BOLETÍN
    # ----------------------------------------------------------
    def _crear(self):
        carpeta = "QUICK ALERT/"
        os.makedirs(carpeta, exist_ok=True)
        nombre = self.var_nombre_archivo.get().strip() or "Modelo Oberon 1.pptx"
        if not nombre.lower().endswith(".pptx"):
            nombre += ".pptx"
        archivo = os.path.abspath(os.path.join(carpeta, nombre))

        contador = 2
        nombre_base, ext = os.path.splitext(nombre)
        while os.path.exists(archivo):
            nombre = f"{nombre_base} ({contador}){ext}"
            archivo = os.path.abspath(os.path.join(carpeta, nombre))
            contador += 1

        self.ultimo_archivo = archivo

        prs = Presentation()
        prs.slide_width = Cm(38.1)
        prs.slide_height = Cm(67.733)
        slide = prs.slides.add_slide(prs.slide_layouts[6])

        slide.shapes.add_picture(self.var_cabezote.get(), Cm(0), Cm(0), Cm(38.1), Cm(17.64))
        add_simple_text(slide, 2, 18.64, 34, 1.56, f"Bogotá, {fecha_formateada}", 30)

        if self.var_icono_boletin.get():
            slide.shapes.add_picture(self.var_icono_boletin.get(), Cm(27.5), Cm(18.64), Cm(1.5), Cm(1.5))
        add_simple_text(slide, 29, 18.64, 7.9, 1.56, self.var_boletin.get(), 30)

        if self.var_icono_ubicacion.get():
            slide.shapes.add_picture(self.var_icono_ubicacion.get(), Cm(2), Cm(20.64), Cm(1.3), Cm(1.3))
        add_simple_text(slide, 3.5, 20.64, 20.05, 1.56, self.var_ubicacion.get(), 30, True)

        if self.var_icono_titulo.get():
            slide.shapes.add_picture(self.var_icono_titulo.get(), Cm(2), Cm(22.64), Cm(1.3), Cm(1.3))
        add_simple_text(slide, 3.5, 22.64, 34.34, 1.37, self.var_titulo.get(), 23)

        txt_p = self.txt_principal.get("1.0", "end-1c")
        add_rich_text_box(slide, 2, 24.64, 34.56, 7.07, txt_p, 23, get_bold_ranges(self.txt_principal))

        slide.shapes.add_picture(self.var_evidencia.get(), Cm(2), Cm(32.64), Cm(16.35), Cm(9.2))

        txt_s = self.txt_secundario.get("1.0", "end-1c")
        add_rich_text_box(slide, 20, 32.64, 16.82, 7.07, txt_s, 23, get_bold_ranges(self.txt_secundario))

        add_simple_text(slide, 1.8, 41.7, 5, 0.79, self.var_fuente.get(), 15)

        add_simple_text(slide, 1.8, 42.7, 5, 0.79, "IMPACTO:", 23, True)

        add_simple_text(slide, 6, 42.7, 10, 0.79, self.var_impacto.get(), 23, False)

        if self.var_icono_recomendaciones.get():
            slide.shapes.add_picture(self.var_icono_recomendaciones.get(), Cm(2), Cm(44), Cm(1.5), Cm(1.5))
        add_simple_text(slide, 3.7, 44, 15.62, 1.56, "RECOMENDACIONES", 27)

        txt_r = self.txt_recomendaciones.get("1.0", "end-1c")
        add_rich_text_box(slide, 2, 45.5, 34.56, 15, txt_r, 23, get_bold_ranges(self.txt_recomendaciones))

        prs.save(archivo)
        nombre_archivo = os.path.splitext(os.path.basename(archivo))[0]
        messagebox.showinfo("Éxito", f"Boletín: {nombre_archivo} creado correctamente.")

# ----------------------------------------------------------
# EJECUCIÓN
# ----------------------------------------------------------
if __name__ == "__main__":
    AppBoletin().mainloop()
